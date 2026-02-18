"""Pre-built PowerPoint helper for consistent, professional slides.

Usage in execute_code scripts:
    import sys; sys.path.insert(0, '/app')
    from grok.pptx_template import Deck

    deck = Deck("Presentation Title")
    deck.add_title_slide("Main Title", "Subtitle or tagline")
    deck.add_section_slide("Section Header", "Brief description")
    deck.add_content_slide("Slide Title", [
        "First key point with supporting detail",
        "Second key point with evidence or data",
        "Third key point with implication",
    ])
    deck.add_two_column_slide("Comparison Title",
        left_title="Option A", left_points=["Point 1", "Point 2"],
        right_title="Option B", right_points=["Point 1", "Point 2"],
    )
    deck.add_quote_slide("Notable quote here", "— Attribution")
    deck.add_closing_slide("Key Takeaway", "Contact or next steps")
    deck.save("/tmp/output/presentation.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# Color palette
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x96, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x3F)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xFA)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_points(text_frame, points, font_size=16, color=WHITE, spacing=Pt(8), font_name="Calibri"):
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0  # No sub-bullets


class Deck:
    def __init__(self, title="Presentation"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.title = title

    def _blank_slide(self):
        layout = self.prs.slide_layouts[6]  # Blank
        return self.prs.slides.add_slide(layout)

    def add_title_slide(self, title, subtitle=""):
        slide = self._blank_slide()
        _set_slide_bg(slide, DARK_BG)

        # Accent bar
        shape = slide.shapes.add_shape(
            1, Inches(0.8), Inches(2.8), Inches(1.5), Pt(4)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()

        _add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(1.5),
                     title, font_size=44, color=WHITE, bold=True)

        if subtitle:
            _add_textbox(slide, Inches(0.8), Inches(4.5), Inches(10), Inches(1),
                         subtitle, font_size=22, color=LIGHT_GRAY)

    def add_section_slide(self, heading, description=""):
        slide = self._blank_slide()
        _set_slide_bg(slide, ACCENT)

        _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                     heading, font_size=40, color=WHITE, bold=True)

        if description:
            _add_textbox(slide, Inches(1), Inches(4.2), Inches(10), Inches(1),
                         description, font_size=20, color=RGBColor(0xE0, 0xF0, 0xFF))

    def add_content_slide(self, title, points, footnote=""):
        slide = self._blank_slide()
        _set_slide_bg(slide, DARK_BG)

        # Title
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                     title, font_size=30, color=ACCENT, bold=True)

        # Accent line under title
        shape = slide.shapes.add_shape(
            1, Inches(0.8), Inches(1.2), Inches(11.5), Pt(2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()

        # Content points
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(11), Inches(4.8)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _add_points(tf, points, font_size=18, color=WHITE)

        if footnote:
            _add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
                         footnote, font_size=12, color=MID_GRAY)

    def add_two_column_slide(self, title, left_title, left_points,
                              right_title, right_points):
        slide = self._blank_slide()
        _set_slide_bg(slide, DARK_BG)

        # Title
        _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                     title, font_size=30, color=ACCENT, bold=True)

        # Left column header
        _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6),
                     left_title, font_size=22, color=WHITE, bold=True)

        # Left points
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _add_points(tf, left_points, font_size=16, color=LIGHT_GRAY)

        # Divider
        shape = slide.shapes.add_shape(
            1, Inches(6.6), Inches(1.5), Pt(2), Inches(4.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = MID_GRAY
        shape.line.fill.background()

        # Right column header
        _add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.6),
                     right_title, font_size=22, color=WHITE, bold=True)

        # Right points
        txBox = slide.shapes.add_textbox(
            Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _add_points(tf, right_points, font_size=16, color=LIGHT_GRAY)

    def add_quote_slide(self, quote, attribution=""):
        slide = self._blank_slide()
        _set_slide_bg(slide, DARK_BG)

        _add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(3),
                     f'"{quote}"', font_size=28, color=WHITE,
                     alignment=PP_ALIGN.CENTER)

        if attribution:
            _add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.6),
                         attribution, font_size=18, color=ACCENT,
                         alignment=PP_ALIGN.CENTER)

    def add_closing_slide(self, headline, subtext=""):
        slide = self._blank_slide()
        _set_slide_bg(slide, DARK_BG)

        _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                     headline, font_size=40, color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        if subtext:
            _add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(1),
                         subtext, font_size=20, color=LIGHT_GRAY,
                         alignment=PP_ALIGN.CENTER)

    def save(self, path="/tmp/output/presentation.pptx"):
        self.prs.save(path)
